import pandas as pd
from langchain_community.document_loaders import UnstructuredPowerPointLoader
import json
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
from unstructured.documents.elements import PageBreak
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
import os

def get_slide_notes(ppt_path):
    ppt=Presentation(ppt_path)
    notes = []
    for page, slide in enumerate(ppt.slides):
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append((page,textNote))
    return notes

def get_slides(file) -> list[int]:
    elements = partition_pptx(file)
    slides = []
    page_number = ""
    file_name = ""
    slide = ""
    for e in elements :
        # print(e.metadata.fields)
        if isinstance(e, PageBreak):
            slides.append({"slide_text":slide, "page_number":page_number, "file_name":file_name})
            slide = ""
        else:
            slide += e.text+" "
            file_name = e.metadata.filename
            page_number = e.metadata.page_number
    
    slides.append({"slide_text":slide, "page_number":page_number, "file_name":file_name})
       
    return slides

def create_chunks(docs, chunk_size=1000, chunk_overlap=100):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    texts = text_splitter.split_documents(docs)
    return texts

def merge_slides_notes(slides_text, slides_notes):
    slides_val = slides_text
    for index,s_notes in enumerate(slides_notes):
        s_note_val = s_notes[1]
        if len(s_note_val) > 1:
            s_val = slides_val[index]['slide_text']
            s_val = s_val+" "+s_note_val
            slides_val[index]['slide_text'] = s_val  
    return slides_val

def create_json(slides_data, document_url):
    data = []
    dd = []
    for s_data in slides_data:
        doc =  Document(page_content=s_data['slide_text'].strip(), metadata={"source": s_data['file_name'],"page_number":s_data['page_number']})
        dd.append(doc)

    slides_docs = create_chunks(dd)
    for slides in slides_docs:
        page_data = {}
        page_data['text'] = slides.page_content
        page_data['document_name'] = slides.metadata['source']
        page_data['page_number'] = slides.metadata['page_number']
        page_data['document_url'] = document_url
        page_data['title'] = f"{slides.metadata['source']} - page_number {slides.metadata['page_number']}"
        data.append(page_data)
    return data

def parse_pptx(ppt_path, document_url, index_name, custom_columns):
    output_data_path = os.getcwd()
    ppt_file_path_arr = ppt_path.split("/")
    ppt_file_name = ppt_file_path_arr[-1]
    ppt_file_name_without_extension = ppt_file_name.rsplit('.', 1)[0]
    
    # json_path = "../data/sample_data.json"
    slides_data = get_slides(ppt_path)
    slides_notes_data = get_slide_notes(ppt_path)
    slides_merg_data = merge_slides_notes(slides_data, slides_notes_data)
    data = create_json(slides_merg_data, document_url)

    # # writing to json file
    # with open(json_path, "w") as final:
    #     json.dump(data, final)

    os.makedirs(os.path.dirname(f"{output_data_path}/open_search_data/output.json"), exist_ok=True)
    with open(f"{output_data_path}/open_search_data/output_{ppt_file_name_without_extension}.json", "w") as f:
        for content in data:
            f.write(json.dumps({"index": {"_index": index_name}}) + '\n')
            f.write(json.dumps(content) + "\n")

    print(f"Processed {ppt_path}")

  